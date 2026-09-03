"""Rendu des aperçus Word, Excel et PowerPoint.

Trois défauts de mise en forme étaient visibles à l'écran :

- Word : ``doc.paragraphs`` et ``doc.tables`` étant deux collections séparées,
  les lire l'une après l'autre rejetait tous les tableaux à la fin ; et les
  puces, rendues en ``<p>``, perdaient toute marque.
- Excel : une feuille d'une seule ligne ouvrait un ``<thead>`` jamais fermé et
  fermait un ``<tbody>`` jamais ouvert.
- PowerPoint : le titre de diapositive était noyé parmi les autres paragraphes.
"""

import io

import pytest
from django.core.files.uploadedfile import SimpleUploadedFile

from apps.documents.models import Document
from apps.documents.services import _format_cell, generate_preview_html

pytestmark = pytest.mark.django_db


def _document(name, payload):
    return Document(
        title=name,
        file=SimpleUploadedFile(name, payload),
        file_name=name,
    )


def _docx_bytes(build):
    from docx import Document as DocxDocument

    doc = DocxDocument()
    build(doc)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _xlsx_bytes(rows, sheet_title='Feuille'):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_title
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------

def test_docx_keeps_tables_where_the_document_puts_them():
    def build(doc):
        doc.add_paragraph('Avant le tableau')
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = 'Colonne A'
        table.cell(1, 0).text = 'Valeur A'
        doc.add_paragraph('Apres le tableau')

    html, error = generate_preview_html(_document('ordre.docx', _docx_bytes(build)))

    assert error is None
    avant = html.index('Avant le tableau')
    tableau = html.index('Colonne A')
    apres = html.index('Apres le tableau')
    # Le tableau était auparavant rejeté après « Apres le tableau ».
    assert avant < tableau < apres


def test_docx_renders_bullets_as_a_list():
    def build(doc):
        doc.add_paragraph('Premier point', style='List Bullet')
        doc.add_paragraph('Second point', style='List Bullet')

    html, error = generate_preview_html(_document('puces.docx', _docx_bytes(build)))

    assert error is None
    assert '<ul>' in html and '</ul>' in html
    assert html.count('<li>') == 2


def test_docx_renders_numbered_items_as_an_ordered_list():
    def build(doc):
        doc.add_paragraph('Etape une', style='List Number')
        doc.add_paragraph('Etape deux', style='List Number')

    html, error = generate_preview_html(_document('num.docx', _docx_bytes(build)))

    assert error is None
    assert '<ol>' in html and html.count('<li>') == 2


def test_docx_closes_the_list_before_a_following_paragraph():
    def build(doc):
        doc.add_paragraph('Un point', style='List Bullet')
        doc.add_paragraph('Un paragraphe ordinaire')

    html, error = generate_preview_html(_document('mix.docx', _docx_bytes(build)))

    assert error is None
    assert html.index('</ul>') < html.index('Un paragraphe ordinaire')


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------

def test_xlsx_single_row_sheet_produces_balanced_markup():
    """Le cas qui cassait le tableau à l'écran."""
    html, error = generate_preview_html(
        _document('une_ligne.xlsx', _xlsx_bytes([('Nom', 'Montant')]))
    )

    assert error is None
    assert html.count('<thead') == html.count('</thead>')
    assert html.count('<tbody') == html.count('</tbody>')
    # Aucun </tbody> orphelin quand il n'y a pas de corps de tableau.
    assert '<tbody>' not in html


def test_xlsx_multi_row_sheet_opens_and_closes_a_body():
    html, error = generate_preview_html(
        _document('lignes.xlsx', _xlsx_bytes([('Nom', 'Montant'), ('Dupont', 120)]))
    )

    assert error is None
    assert html.count('<thead') == html.count('</thead>')
    assert html.count('<tbody>') == html.count('</tbody>') == 1


def test_xlsx_skips_a_sheet_with_no_content():
    html, error = generate_preview_html(
        _document('vide.xlsx', _xlsx_bytes([], sheet_title='Rien'))
    )

    assert error is None
    # La feuille vide affichait un titre suivi d'un tableau sans lignes.
    assert 'Rien' not in html
    assert 'Classeur vide' in html


def test_cell_values_are_formatted_for_reading():
    from datetime import date, datetime

    assert _format_cell(date(2026, 1, 5)) == '05/01/2026'
    assert _format_cell(datetime(2026, 1, 5, 0, 0)) == '05/01/2026'
    assert _format_cell(datetime(2026, 1, 5, 14, 30)) == '05/01/2026 14:30'
    # 120.0 s'affichait tel quel pour un montant saisi sans décimale.
    assert _format_cell(120.0) == '120'
    assert _format_cell(120.5) == '120.5'
    assert _format_cell(None) == ''


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def _pptx_bytes(title_text, body_text):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = body_text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_marks_the_slide_title():
    html, error = generate_preview_html(
        _document('deck.pptx', _pptx_bytes('Le titre', 'Le corps'))
    )

    assert error is None
    # Le titre était rendu comme n'importe quel autre paragraphe.
    assert '<p class="fw-bold mb-2">Le titre</p>' in html
    assert '<p>Le corps</p>' in html
    assert 'Diapositive 1' in html
