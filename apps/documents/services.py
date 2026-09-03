import mimetypes
import logging
from datetime import timedelta
from django.conf import settings
from django.core.mail import EmailMultiAlternatives
from django.db.models import Count, Sum, Q
from django.template.loader import render_to_string
from django.urls import reverse
from django.utils import timezone
from django.utils.html import strip_tags

from .models import Document, DocumentCategory, DocumentShare, DocumentAccess

logger = logging.getLogger(__name__)

# ─── Validation ────────────────────────────────────────────────

ALLOWED_EXTENSIONS = {
    'document': {'pdf', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt', 'txt', 'csv', 'zip', 'odt', 'ods', 'rtf'},
    'image': {'jpg', 'jpeg', 'png', 'gif', 'webp', 'svg', 'bmp'},
    'audio': {'mp3', 'wav', 'ogg', 'm4a', 'flac', 'aac', 'wma'},
    'video': {'mp4', 'webm', 'mov', 'mkv', 'avi', 'wmv', 'm4v'},
}

BLOCKED_EXTENSIONS = {'exe', 'bat', 'sh', 'py', 'js', 'cmd', 'com', 'msi', 'scr', 'vbs', 'ps1', 'dll'}

MAX_FILE_SIZES = {
    'document': 20 * 1024 * 1024,       # 20 Mo
    'image': 10 * 1024 * 1024,           # 10 Mo
    'audio': 100 * 1024 * 1024,          # 100 Mo
    'video': 500 * 1024 * 1024,          # 500 Mo
}

ALL_ALLOWED = set()
for exts in ALLOWED_EXTENSIONS.values():
    ALL_ALLOWED |= exts


def detect_media_type(filename, mime_type=''):
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    if ext in BLOCKED_EXTENSIONS:
        return None
    for media_type, extensions in ALLOWED_EXTENSIONS.items():
        if ext in extensions:
            return media_type
    if mime_type:
        if mime_type.startswith('image/'):
            return 'image'
        if mime_type.startswith('audio/'):
            return 'audio'
        if mime_type.startswith('video/'):
            return 'video'
    return None


def validate_file(uploaded_file):
    """Valide un fichier uploadé. Retourne (media_type, error_message)."""
    ext = uploaded_file.name.rsplit('.', 1)[-1].lower() if '.' in uploaded_file.name else ''

    if ext in BLOCKED_EXTENSIONS:
        return None, f"Type de fichier interdit (.{ext})"

    mime_type = uploaded_file.content_type or ''
    media_type = detect_media_type(uploaded_file.name, mime_type)

    if not media_type:
        return None, f"Type de fichier non supporté (.{ext})"

    max_size = MAX_FILE_SIZES.get(media_type, 20 * 1024 * 1024)
    if uploaded_file.size > max_size:
        max_display = max_size // (1024 * 1024)
        return None, f"Fichier trop volumineux ({uploaded_file.size // (1024 * 1024)} Mo). Maximum : {max_display} Mo pour les {media_type}s"

    return media_type, None


def get_mime_type(filename):
    mime, _ = mimetypes.guess_type(filename)
    return mime or 'application/octet-stream'


# ─── Statistiques ──────────────────────────────────────────────

def get_documents_stats(user=None):
    qs = Document.accessible_queryset(user) if user else Document.objects.all()

    total = qs.count()
    total_size = qs.aggregate(s=Sum('file_size'))['s'] or 0
    by_media = dict(qs.values_list('media_type').annotate(c=Count('id')).values_list('media_type', 'c'))
    by_category = dict(
        qs.filter(category__isnull=False)
        .values_list('category__name')
        .annotate(c=Count('id'))
        .values_list('category__name', 'c')
    )
    confidential = qs.filter(is_confidential=True).count()
    recent_cutoff = timezone.now() - timedelta(days=30)
    recent_count = qs.filter(created_at__gte=recent_cutoff).count()
    recent = qs.order_by('-created_at')[:5]

    return {
        'total': total,
        'total_size': total_size,
        'total_size_display': _format_size(total_size),
        'by_media': by_media,
        'by_category': by_category,
        'confidential': confidential,
        'recent_count': recent_count,
        'recent': recent,
    }


def _format_size(size):
    if size < 1024:
        return f"{size} o"
    elif size < 1024 * 1024:
        return f"{size / 1024:.1f} Ko"
    elif size < 1024 * 1024 * 1024:
        return f"{size / (1024 * 1024):.1f} Mo"
    return f"{size / (1024 * 1024 * 1024):.2f} Go"


# ─── Partage par email ────────────────────────────────────────

def share_document_by_email(document, user, recipient_email, recipient_name='', message='', request=None):
    """Crée un partage et envoie l'email. Retourne le DocumentShare."""
    share = DocumentShare.objects.create(
        document=document,
        shared_by=user,
        recipient_email=recipient_email,
        recipient_name=recipient_name,
        message=message,
        expires_at=timezone.now() + timedelta(days=7),
    )

    # Log l'accès
    DocumentAccess.objects.create(
        document=document,
        user=user,
        action='share',
        ip_address=_get_client_ip(request) if request else None,
    )

    # Construire le lien
    if request:
        share_url = request.build_absolute_uri(
            reverse('documents:shared_access', kwargs={'token': share.share_token})
        )
    else:
        share_url = reverse('documents:shared_access', kwargs={'token': share.share_token})

    # Envoyer l'email
    sender_name = user.get_full_name() or user.username
    context = {
        'document': document,
        'share': share,
        'share_url': share_url,
        'shared_by': sender_name,
        'sender_name': sender_name,
        'recipient_name': recipient_name,
        'message': message,
        'document_title': document.title,
        'file_size': document.file_size_display,
        'media_type': document.get_media_type_display(),
        'expires_at': share.expires_at,
    }

    html_content = render_to_string('documents/emails/document_shared.html', context)
    text_content = strip_tags(html_content)

    subject = f"Document partagé : {document.title}"
    from_email = settings.DEFAULT_FROM_EMAIL

    email_obj = EmailMultiAlternatives(subject, text_content, from_email, [recipient_email])
    email_obj.attach_alternative(html_content, 'text/html')

    email_sent = False
    email_error = ''
    try:
        email_obj.send()
        email_sent = True
        logger.info(f"Email de partage envoyé à {recipient_email} pour le document {document.id}")
    except Exception as e:
        email_error = str(e)
        logger.error(f"Erreur envoi email partage: {e}")

    # Log dans EmailLog
    try:
        from apps.communication.models import EmailLog
        EmailLog.objects.create(
            recipient_email=recipient_email,
            recipient_name=recipient_name,
            subject=subject,
            body=text_content[:2000],
            status='sent' if email_sent else 'failed',
            sent_at=timezone.now() if email_sent else None,
            error_message=email_error,
        )
    except Exception as e:
        logger.warning(f"Impossible de créer EmailLog: {e}")

    share._email_sent = email_sent
    share._email_error = email_error
    return share


def log_access(document, user, action, request=None):
    DocumentAccess.objects.create(
        document=document,
        user=user,
        action=action,
        ip_address=_get_client_ip(request) if request else None,
    )


def _get_client_ip(request):
    if not request:
        return None
    x_forwarded = request.META.get('HTTP_X_FORWARDED_FOR')
    if x_forwarded:
        return x_forwarded.split(',')[0].strip()
    return request.META.get('REMOTE_ADDR')


# ─── Prévisualisation de contenu ──────────────────────────────

TEXT_EXTENSIONS = {'txt', 'csv', 'json', 'xml', 'md', 'log', 'yml', 'yaml', 'ini', 'conf', 'html', 'css'}
MAX_PREVIEW_SIZE = 512 * 1024  # 512 Ko max pour les previews texte
# Un classeur est charge entierement en memoire avant lecture : plafonner
# evite qu'un fichier de plusieurs dizaines de Mo passe par le processus web.
MAX_OFFICE_PREVIEW_SIZE = 20 * 1024 * 1024  # 20 Mo


def generate_preview_html(document):
    """Génère un aperçu HTML du contenu d'un fichier.
    Retourne (html_content, error_message)."""
    ext = document.extension

    try:
        if ext in TEXT_EXTENSIONS:
            return _preview_text(document, ext)
        elif ext == 'docx':
            return _preview_docx(document)
        elif ext == 'xlsx':
            return _preview_xlsx(document)
        elif ext == 'pptx':
            return _preview_pptx(document)
        elif ext in {'doc', 'xls', 'ppt'}:
            return _preview_legacy_office(document, ext)
    except Exception as e:
        logger.error(f"Erreur preview document {document.id}: {e}")
        return None, f"Impossible de générer l'aperçu : {e}"

    return None, "Format non supporté pour la prévisualisation."


def _preview_text(document, ext):
    """Aperçu pour fichiers texte (txt, csv, json, xml, md, etc.)."""
    import markdown as md_lib
    from django.utils.html import escape

    file_obj = document.file.open('rb')
    try:
        raw = file_obj.read(MAX_PREVIEW_SIZE)
    finally:
        file_obj.close()

    # Détecter l'encodage
    for encoding in ('utf-8', 'latin-1', 'cp1252'):
        try:
            content = raw.decode(encoding)
            break
        except (UnicodeDecodeError, ValueError):
            continue
    else:
        return None, "Impossible de décoder le fichier texte."

    truncated = len(raw) >= MAX_PREVIEW_SIZE

    if ext == 'md':
        html = md_lib.markdown(content, extensions=['tables', 'fenced_code', 'nl2br'])
        return _wrap_preview(html, truncated, 'Markdown'), None

    if ext == 'csv':
        return _preview_csv(content, truncated), None

    # Texte brut, json, xml, yaml, etc.
    label = ext.upper()
    escaped = escape(content)
    html = f'<pre class="preview-code">{escaped}</pre>'
    return _wrap_preview(html, truncated, label), None


def _preview_csv(content, truncated):
    """Rend un CSV en tableau HTML (max 200 lignes)."""
    import csv
    import io
    from django.utils.html import escape

    reader = csv.reader(io.StringIO(content))
    rows = []
    for i, row in enumerate(reader):
        if i >= 200:
            truncated = True
            break
        rows.append(row)

    if not rows:
        return _wrap_preview('<p class="text-muted">Fichier CSV vide.</p>', False, 'CSV')

    html = '<div class="table-responsive"><table class="table table-sm table-bordered table-striped mb-0">'
    # Première ligne = entête
    html += '<thead class="table-dark"><tr>'
    for cell in rows[0]:
        html += f'<th>{escape(cell)}</th>'
    html += '</tr></thead><tbody>'
    for row in rows[1:]:
        html += '<tr>'
        for cell in row:
            html += f'<td>{escape(cell)}</td>'
        html += '</tr>'
    html += '</tbody></table></div>'
    return _wrap_preview(html, truncated, 'CSV')


def _iter_docx_blocks(doc):
    """Parcourt paragraphes et tableaux dans l'ordre du document.

    python-docx expose doc.paragraphs et doc.tables comme deux collections
    separees. Les lire l'une apres l'autre renvoyait tous les tableaux a la
    fin de l'apercu : un document « intro, tableau, conclusion » s'affichait
    « intro, conclusion, tableau ». Seul le corps XML porte l'ordre reel.
    """
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in doc.element.body.iterchildren():
        if child.tag == qn('w:p'):
            yield Paragraph(child, doc)
        elif child.tag == qn('w:tbl'):
            yield Table(child, doc)


def _docx_table_html(table, escape):
    parts = ['<div class="table-responsive"><table class="table table-sm table-bordered mb-3">']
    for i, row in enumerate(table.rows):
        tag = 'th' if i == 0 else 'td'
        if i == 0:
            parts.append('<thead>')
        elif i == 1:
            parts.append('</thead><tbody>')
        parts.append('<tr>')
        for cell in row.cells:
            parts.append('<%s>%s</%s>' % (tag, escape(cell.text), tag))
        parts.append('</tr>')
    # Fermer ce qui a reellement ete ouvert : un tableau d'une seule ligne
    # n'a jamais de <tbody>.
    parts.append('</thead>' if len(table.rows) <= 1 else '</tbody>')
    parts.append('</table></div>')
    return parts


def _preview_docx(document):
    """Apercu d'un fichier Word (.docx) via python-docx."""
    from docx import Document as DocxDocument
    from django.utils.html import escape

    file_obj = document.file.open('rb')
    try:
        doc = DocxDocument(file_obj)
    finally:
        file_obj.close()

    parts = []
    state = {'list': None}

    def close_list():
        if state['list']:
            parts.append('</%s>' % state['list'])
            state['list'] = None

    for block in _iter_docx_blocks(doc):
        if hasattr(block, 'rows'):
            close_list()
            parts.extend(_docx_table_html(block, escape))
            continue

        text = block.text.strip()
        if not text:
            continue

        style_name = (block.style.name or '').lower()

        # Les puces et numeros de Word portent un style de liste. Rendus en
        # <p>, ils perdaient toute marque : une enumeration devenait une suite
        # de phrases sans lien apparent.
        if 'list' in style_name:
            wanted = 'ol' if 'number' in style_name else 'ul'
            if state['list'] != wanted:
                close_list()
                parts.append('<%s>' % wanted)
                state['list'] = wanted
            parts.append('<li>%s</li>' % escape(text))
            continue

        close_list()

        if 'heading 1' in style_name:
            parts.append('<h3>%s</h3>' % escape(text))
        elif 'heading 2' in style_name:
            parts.append('<h4>%s</h4>' % escape(text))
        elif 'heading 3' in style_name:
            parts.append('<h5>%s</h5>' % escape(text))
        elif 'title' in style_name:
            parts.append('<h2>%s</h2>' % escape(text))
        elif 'quote' in style_name:
            parts.append('<blockquote class="blockquote">%s</blockquote>' % escape(text))
        else:
            parts.append('<p>%s</p>' % escape(text))

    close_list()

    if not parts:
        return _wrap_preview('<p class="text-muted">Document vide.</p>', False, 'Word'), None

    html = '\n'.join(parts)
    return _wrap_preview(html, False, 'Word (.docx)'), None


def _format_cell(value):
    """Rend une valeur de cellule lisible plutot que sa repr Python."""
    from datetime import date, datetime, time

    if value is None:
        return ''
    if isinstance(value, datetime):
        # str() donnait « 2026-01-05 00:00:00 » sur une simple date.
        if (value.hour, value.minute, value.second) == (0, 0, 0):
            return value.strftime('%d/%m/%Y')
        return value.strftime('%d/%m/%Y %H:%M')
    if isinstance(value, date):
        return value.strftime('%d/%m/%Y')
    if isinstance(value, time):
        return value.strftime('%H:%M')
    if isinstance(value, float) and value.is_integer():
        # 120.0 pour un montant entier saisi sans decimale.
        return str(int(value))
    return str(value)


def _preview_xlsx(document):
    """Apercu d'un fichier Excel (.xlsx) via openpyxl."""
    import io as _io

    from openpyxl import load_workbook
    from django.utils.html import escape

    # read_only=True lit la feuille en flux : openpyxl a encore besoin du
    # fichier au moment d'itérer les lignes. Le fermer juste apres
    # load_workbook faisait echouer chaque apercu Excel sur « I/O operation on
    # closed file ». On charge donc les octets en memoire d'abord.
    file_obj = document.file.open('rb')
    try:
        payload = file_obj.read(MAX_OFFICE_PREVIEW_SIZE)
    finally:
        file_obj.close()

    wb = load_workbook(_io.BytesIO(payload), read_only=True, data_only=True)

    parts = []
    for sheet_name in wb.sheetnames[:5]:  # Max 5 feuilles
        ws = wb[sheet_name]

        rows = []
        for row in ws.iter_rows(max_row=100, values_only=True):
            if any(cell is not None for cell in row):
                rows.append(row)

        # Une feuille vide affichait un titre suivi d'un tableau sans contenu.
        if not rows:
            continue

        parts.append('<h5 class="mt-3 mb-2"><i class="bi bi-table me-1"></i>%s</h5>' % escape(sheet_name))
        parts.append('<div class="table-responsive"><table class="table table-sm table-bordered table-striped mb-3">')
        parts.append('<thead class="table-dark"><tr>')
        for cell in rows[0]:
            parts.append('<th>%s</th>' % escape(_format_cell(cell)))
        parts.append('</tr></thead>')

        # Le <tbody> n'etait ouvert qu'a partir de la deuxieme ligne, alors
        # qu'il etait toujours ferme en sortie de boucle : une feuille d'une
        # seule ligne produisait un <thead> ouvert et un </tbody> orphelin.
        if len(rows) > 1:
            parts.append('<tbody>')
            for row in rows[1:]:
                parts.append('<tr>')
                for cell in row:
                    parts.append('<td>%s</td>' % escape(_format_cell(cell)))
                parts.append('</tr>')
            parts.append('</tbody>')

        parts.append('</table></div>')

        if len(rows) >= 100:
            parts.append(
                '<p class="text-muted"><em>Apercu limite a 100 lignes pour « %s ».</em></p>'
                % escape(sheet_name)
            )

    wb.close()

    if not parts:
        return _wrap_preview('<p class="text-muted">Classeur vide.</p>', False, 'Excel'), None

    html = '\n'.join(parts)
    return _wrap_preview(html, False, 'Excel (.xlsx)'), None


def _preview_pptx(document):
    """Apercu d'un fichier PowerPoint (.pptx) via python-pptx."""
    from pptx import Presentation
    from django.utils.html import escape

    file_obj = document.file.open('rb')
    try:
        prs = Presentation(file_obj)
    finally:
        file_obj.close()

    parts = []
    for i, slide in enumerate(prs.slides):
        title = ''
        body = []
        title_shape = slide.shapes.title

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # Le titre de la diapositive se retrouvait noye parmi les
                # autres paragraphes, tous rendus en <p>.
                if not title and title_shape is not None and shape == title_shape:
                    title = text
                else:
                    body.append(text)

        parts.append('<div class="preview-slide"><h5 class="mb-2">'
                     '<i class="bi bi-easel me-1"></i>Diapositive %d</h5>' % (i + 1))
        if title:
            parts.append('<p class="fw-bold mb-2">%s</p>' % escape(title))
        for text in body:
            parts.append('<p>%s</p>' % escape(text))
        if not title and not body:
            parts.append('<p class="text-muted mb-0"><em>Diapositive sans texte.</em></p>')
        parts.append('</div>')

    if not parts:
        return _wrap_preview('<p class="text-muted">Presentation vide.</p>', False, 'PowerPoint'), None

    html = '\n'.join(parts)
    return _wrap_preview(html, False, 'PowerPoint (.pptx)'), None


def _preview_legacy_office(document, ext):
    labels = {
        'doc': 'Word 97-2003 (.doc)',
        'xls': 'Excel 97-2003 (.xls)',
        'ppt': 'PowerPoint 97-2003 (.ppt)',
    }
    label = labels.get(ext, ext.upper())
    html = (
        '<div class="alert alert-warning mb-0">'
        '<i class="bi bi-exclamation-triangle me-1"></i>'
        f"L'aperçu intégré ne peut pas lire directement les anciens fichiers {label}. "
        'Téléchargez le fichier ou enregistrez-le au format moderne '
        '(.docx, .xlsx ou .pptx) pour obtenir une prévisualisation dans l\'application.'
        '</div>'
    )
    return _wrap_preview(html, False, label), None


def _wrap_preview(inner_html, truncated, format_label):
    """Enveloppe le contenu prévisualisé avec un cadre."""
    badge = f'<span class="badge bg-secondary me-2">{format_label}</span>'
    warning = ''
    if truncated:
        warning = '<div class="alert alert-info py-1 px-2 mt-2 mb-0" style="font-size: 0.8rem;"><i class="bi bi-info-circle me-1"></i>Aperçu tronqué — téléchargez pour voir l\'intégralité.</div>'
    return f'{badge}{warning}<div class="preview-body mt-2">{inner_html}</div>'


# ─── Catégories par défaut ─────────────────────────────────────

DEFAULT_CATEGORIES = [
    {'name': 'Procès-verbaux', 'icon': 'bi-journal-text', 'color': '#366092', 'order': 1},
    {'name': 'Statuts & Règlements', 'icon': 'bi-file-earmark-ruled', 'color': '#6c757d', 'order': 2},
    {'name': 'Reçus fiscaux', 'icon': 'bi-receipt', 'color': '#28a745', 'order': 3},
    {'name': 'Rapports financiers', 'icon': 'bi-graph-up', 'color': '#dc3545', 'order': 4},
    {'name': 'Courriers officiels', 'icon': 'bi-envelope-paper', 'color': '#0d6efd', 'order': 5},
    {'name': 'Photos & Médias', 'icon': 'bi-camera', 'color': '#e83e8c', 'order': 6},
    {'name': 'Prédications', 'icon': 'bi-mic', 'color': '#7c3aed', 'order': 7},
    {'name': 'Louange & Worship', 'icon': 'bi-music-note-beamed', 'color': '#f59e0b', 'order': 8},
    {'name': 'Formulaires', 'icon': 'bi-ui-checks-grid', 'color': '#fd7e14', 'order': 9},
    {'name': 'Autres', 'icon': 'bi-folder2-open', 'color': '#6f42c1', 'order': 10},
]


def create_default_categories():
    created = 0
    for cat_data in DEFAULT_CATEGORIES:
        _, is_new = DocumentCategory.objects.get_or_create(
            name=cat_data['name'],
            defaults=cat_data,
        )
        if is_new:
            created += 1
    return created
